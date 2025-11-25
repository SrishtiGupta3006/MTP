# build_runtime_enforcement_ppt.py
# Creates: Compositional_Schemes_Runtime_Enforcement_Final.pptx
# Requirements: pip install python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# === SETTINGS ===
TITLE_FNT = Pt(28)
BODY_FNT = Pt(20)
CAPTION_FNT = Pt(16)
BLUE = RGBColor(0, 51, 102)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(80, 80, 80)

LOGO_FILE = "iit-bhubaneswar-logo.jpg"
IMG_STRICT = "strict enforcer.png"
IMG_LE_MONO = "least_effort_monolithic.png"
IMG_LE_PAR1 = "least_effort_parallel_1.png"
IMG_LE_PAR2 = "least_effort_parallel_2.png"
IMG_EXCL = "exclusive_modified_automaon.png"

# === PRESENTATION SETUP ===
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
layout_blank = prs.slide_layouts[6]

def add_title_slide():
    slide = prs.slides.add_slide(layout_blank)
    # Title
    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(10), Inches(1.2))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Compositional Schemes in Runtime Enforcement"
    p.font.bold = True
    p.font.size = TITLE_FNT
    p.font.color.rgb = BLUE

    # Subtitle info
    meta = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(10), Inches(2))
    mtf = meta.text_frame
    for line in [
        "Mid-Term Project Presentation",
        "",
        "Srishti Gupta (24CS06016)",
        "Under the supervision of Dr. Srinivas Pinisetty",
        "School of Electrical and Computer Sciences",
        "Indian Institute of Technology Bhubaneswar",
    ]:
        para = mtf.add_paragraph()
        para.text = line
        para.font.size = BODY_FNT
        para.font.color.rgb = BLACK

    # Logo
    try:
        slide.shapes.add_picture(LOGO_FILE, Inches(10.8), Inches(0.6), height=Inches(1.3))
    except:
        pass

def add_slide(title, bullets):
    slide = prs.slides.add_slide(layout_blank)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = TITLE_FNT
    p.font.color.rgb = BLUE

    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    for line in bullets:
        para = tf.add_paragraph()
        para.text = line
        para.font.size = BODY_FNT
        para.font.color.rgb = BLACK
        para.space_before = Pt(6)

def add_image_slide(title, caption, img_path, height=4.5):
    slide = prs.slides.add_slide(layout_blank)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = TITLE_FNT
    p.font.color.rgb = BLUE

    try:
        pic = slide.shapes.add_picture(img_path, Inches(2), Inches(1.5), height=Inches(height))
        cap_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(1))
        cap_tf = cap_box.text_frame
        cap = cap_tf.add_paragraph()
        cap.text = caption
        cap.font.size = CAPTION_FNT
        cap.font.color.rgb = GRAY
        cap.alignment = PP_ALIGN.CENTER
    except Exception as e:
        print(f"Could not add image {img_path}: {e}")

# === SLIDES CONTENT ===
add_title_slide()

add_slide("Motivation", [
    "Safety-critical systems must satisfy multiple correctness properties.",
    "Static verification cannot guarantee compliance during execution.",
    "Runtime Enforcement (RE) dynamically monitors and edits system outputs.",
    "Conflicts arise while enforcing multiple properties concurrently.",
    "Need for compositional schemes balancing correctness and permissiveness."
])

add_slide("Problem Statement", [
    "Existing RE mechanisms (Serial, Monolithic) are overly restrictive.",
    "Conflicts between multiple properties lead to inconsistent outputs.",
    "Need modular, scalable, and minimally intrusive compositional RE."
])

# Literature Review (table simulated as bullet summary for simplicity)
add_slide("Literature Review", [
    "• Runtime Enforcement Monitors: Composition, Synthesis and Enforcer (2012): foundational models; lacks conflict resolution.",
    "• Predictive Runtime Enforcement (2015): reduced latency; limited to single property.",
    "• Bounded Memory RE (2019): supports constrained systems; accuracy trade-off.",
    "• Probabilistic Bounded-Memory RE (2024): balances correctness vs performance.",
    "• Compositional Enforcement of RE (2020): serial/parallel composition; open conflicts.",
    "• Compositional RE Revisited (2022): defines sound parallel rules; complex trade-offs."
])

add_slide("Contribution", [
    "Implemented Serial, Monolithic, Least Effort, and Exclusive Enforcers.",
    "Introduced Grouping Techniques: AND, OR, k-of-n, Priority.",
    "Formalized DFA-based runtime enforcement framework.",
    "Analyzed trade-offs between correctness and permissiveness."
])

add_slide("Runtime Enforcement Framework", [
    "Definition: Ensures outputs satisfy property ϕ by modifying events.",
    "Formal Function: Eϕ : Σ* → Σ*",
    "Key Properties: Soundness, Transparency, Monotonicity.",
    "Operation: Monitor → Evaluate → Edit/Release → Output → Repeat."
])

add_slide("Example of Runtime Enforcement", [
    "Property ϕ: Every 'a' must be followed by 'b'.",
    "Input: a c b",
    "Step 1: 'a' blocked (waiting for 'b'); Step 2: 'c' ignored; Step 3: 'b' forms valid prefix.",
    "Result: Output sequence satisfies ϕ maintaining transparency."
])

add_slide("Grouping Techniques", [
    "OR Grouping: system valid if any property holds.",
    "k-of-n Grouping: valid if at least k properties hold.",
    "Priority Grouping: high-priority properties enforced first.",
    "Applications: fault tolerance, redundancy, safety prioritization."
])

add_slide("AND and OR Grouping", [
    "AND: valid only if all properties hold → strongest correctness.",
    "OR: valid if any property holds → maximum permissiveness.",
    "Examples: Safety constraints (AND); backup restoration (OR)."
])

add_slide("Strict Enforcer Types", [
    "Strict Serial: applies properties sequentially (P1→P2→P3).",
    "Strict Monolithic: combined product automaton enforcing all at once.",
    "Strict Parallel: concurrent enforcement; intersection of accepted events.",
    "Trade-off: highest correctness vs lowest permissiveness."
])

add_image_slide("Strict Enforcer Output",
    "Figure 1: Strict Serial Enforcer Output",
    IMG_STRICT)

add_slide("Least Effort Enforcer Types", [
    "Least Effort Serial: sequential, releases event if any enforcer allows it.",
    "Least Effort Monolithic: disjunctive product DFA (A∨); event allowed if ≥1 accepts.",
    "Least Effort Parallel: merges outputs via logical OR; maximizes permissiveness."
])

# Slide 12 – three images
for idx, (img, caption) in enumerate([
    (IMG_LE_MONO, "Figure 2: Least Effort Monolithic Enforcer Output"),
    (IMG_LE_PAR1, "Figure 3: Least Effort Parallel Enforcer Output (1)"),
    (IMG_LE_PAR2, "Figure 4: Least Effort Parallel Enforcer Output (2)"),
]):
    add_image_slide("Least Effort Enforcer Outputs", caption, img, height=3.8)

add_image_slide("Exclusive Enforcer",
    "Figure 5: Modified Automata for Exclusive Enforcer",
    IMG_EXCL)

add_slide("Conclusion", [
    "Explored compositional runtime enforcement schemes.",
    "Implemented Serial, Monolithic, Least Effort, and Exclusive enforcers.",
    "Introduced AND, OR, k-of-n, and Priority groupings.",
    "Established balance among soundness, transparency, and scalability."
])

add_slide("Future Work", [
    "Extend Exclusive Enforcer with adaptive conflict resolution.",
    "Explore k-of-n grouping for configurable tolerance.",
    "Implement Priority-based grouping for safety-critical systems.",
    "Integrate probabilistic and performance-aware composition."
])

add_slide("References", [
    "1. Pinisetty et al., Runtime Enforcement Monitors..., 2012.",
    "2. Falcone et al., Predictive Runtime Enforcement, 2015.",
    "3. Pinisetty et al., Bounded Memory RE, 2019.",
    "4. Shankar et al., Probabilistic Bounded-Memory RE, 2024.",
    "5. Pinisetty et al., Compositional Enforcement of RE, 2020.",
    "6. Pinisetty et al., Compositional RE Revisited, 2022."
])

prs.save("Compositional_Schemes_Runtime_Enforcement_Final.pptx")
print("✅ Presentation generated: Compositional_Schemes_Runtime_Enforcement_Final.pptx")
